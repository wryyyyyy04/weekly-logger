import os
from datetime import date, timedelta

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import fitz

from models import SessionLocal, DailyTodo, ExperimentLog, WeeklySection, Milestone

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 中文字体
FONT_NAME = "微软雅黑"
FONT_NAME_BOLD = "微软雅黑"

SECTION_TITLES = {
    "plan": "一、上周拟定的本周计划：",
    "progress": "二、本周详细研究进展：",
    "pending": "三、后续待解决的问题或研究方案调整：",
    "next_plan": "四、下周详细计划：",
    "assessment": "五、本人本周各课题进展自评：",
}


def _week_range(week_start):
    end = week_start + timedelta(days=6)
    return week_start, end


def _get_week_data(week_start):
    ws, we = _week_range(week_start)
    with SessionLocal() as db:
        todos = db.query(DailyTodo).filter(
            DailyTodo.date >= ws, DailyTodo.date <= we
        ).order_by(DailyTodo.date, DailyTodo.created_at).all()
        experiments = db.query(ExperimentLog).filter(
            ExperimentLog.date >= ws, ExperimentLog.date <= we
        ).order_by(ExperimentLog.date, ExperimentLog.created_at).all()
        sections = db.query(WeeklySection).filter(
            WeeklySection.week_start == ws
        ).all()
        milestones = db.query(Milestone).order_by(Milestone.sort_order).all()
    return todos, experiments, {s.section_type: s.content for s in sections}, milestones


def _add_slide_content(prs, title_text, body_text, title_size=22, body_size=16):
    """Add a slide with title and body content."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = title_text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.name = FONT_NAME_BOLD
        run.font.size = Pt(title_size)
        run.font.bold = True

    if body_text:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in body_text.strip().split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.space_after = Pt(4)
            for run in p.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(body_size)


def generate_pptx(week_start):
    """生成PPTX文件"""
    todos, experiments, sections, milestones = _get_week_data(week_start)
    ws_str = week_start.isoformat()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 封面 ──
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "周报"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.name = FONT_NAME_BOLD
        run.font.size = Pt(44)
        run.font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = f"方予之\n{ws_str}"
    for p in subtitle.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(24)

    # ── 内容页 ──
    # 一、上周拟定的本周计划
    plan_text = sections.get("plan", "")
    if not plan_text:
        # 如果没有手动填写，尝试从上周的next_plan获取
        prev_week = week_start - timedelta(days=7)
        with SessionLocal() as db:
            prev_section = db.query(WeeklySection).filter(
                WeeklySection.week_start == prev_week,
                WeeklySection.section_type == "next_plan"
            ).first()
        if prev_section and prev_section.content:
            plan_text = prev_section.content
    _add_slide_content(prs, SECTION_TITLES["plan"], plan_text)

    # 二、本周详细研究进展
    progress_text = sections.get("progress", "")
    if not progress_text:
        # 自动汇总本周的待办和实验记录
        lines = []
        for e in experiments:
            lines.append(f"{e.date.isoformat()} {e.title}")
            if e.description:
                lines.append(f"  {e.description}")
        for t in todos:
            prefix = "✓" if t.status == "completed" else "○"
            lines.append(f"{t.date.isoformat()} {prefix} {t.content}")
        if not lines:
            lines.append("（本周暂无记录）")
        progress_text = "\n".join(lines)
    _add_slide_content(prs, SECTION_TITLES["progress"], progress_text)

    # 三、后续待解决问题
    _add_slide_content(prs, SECTION_TITLES["pending"],
                       sections.get("pending", ""),
                       body_size=16)

    # 四、下周详细计划
    _add_slide_content(prs, SECTION_TITLES["next_plan"],
                       sections.get("next_plan", ""),
                       body_size=16)

    # 五、课题进展自评
    _add_slide_content(prs, "五、本人本周各课题进展自评（好、中、差）：",
                       sections.get("assessment", ""),
                       body_size=16)

    # 六、成果发表时间节点
    timeline_lines = []
    for m in milestones:
        planned = m.planned_date.isoformat() if m.planned_date else "（待定）"
        actual = m.actual_date.isoformat() if m.actual_date else "（未完成）"
        timeline_lines.append(
            f"{m.name}\n  计划日期：{planned}  实际日期：{actual}\n  {m.notes or ''}"
        )
    _add_slide_content(prs, "六、本人成果发表时间节点更新：",
                       "\n\n".join(timeline_lines),
                       body_size=14)

    # 保存
    pptx_path = os.path.join(OUTPUT_DIR, f"周报-方予之-{ws_str}.pptx")
    prs.save(pptx_path)
    return pptx_path


def generate_pdf(week_start):
    """直接生成PDF文件（用fitz）"""
    todos, experiments, sections, milestones = _get_week_data(week_start)
    ws_str = week_start.isoformat()

    doc = fitz.open()
    page_w, page_h = 595, 842  # A4

    def add_page(text, title=None):
        page = doc.new_page(width=page_w, height=page_h)
        y = 60

        if title:
            rect = fitz.Rect(50, y, page_w - 50, y + 30)
            page.insert_textbox(rect, title, fontsize=16, fontname="china-s",
                                color=(0, 0, 0))
            y += 40

        if text:
            lines = text.strip().split("\n")
            for line in lines:
                if y > page_h - 60:
                    page = doc.new_page(width=page_w, height=page_h)
                    y = 60
                rect = fitz.Rect(60, y, page_w - 60, y + 20)
                page.insert_textbox(rect, line, fontsize=12, fontname="china-s",
                                    color=(0.2, 0.2, 0.2))
                y += 22

        return page

    # 封面
    cover = doc.new_page(width=page_w, height=page_h)
    rect = fitz.Rect(50, 300, page_w - 50, 360)
    cover.insert_textbox(rect, "周报", fontsize=36, fontname="china-s",
                         color=(0, 0, 0), align=1)
    rect2 = fitz.Rect(50, 380, page_w - 50, 440)
    cover.insert_textbox(rect2, f"方予之\n{ws_str}", fontsize=18, fontname="china-s",
                         color=(0.3, 0.3, 0.3), align=1)

    # 各板块
    plan_text = sections.get("plan", "")
    add_page(plan_text, SECTION_TITLES["plan"])

    progress_text = sections.get("progress", "")
    if not progress_text:
        lines = []
        for e in experiments:
            lines.append(f"{e.date.isoformat()} {e.title}")
            if e.description:
                lines.append(f"  {e.description}")
        for t in todos:
            prefix = "✓" if t.status == "completed" else "○"
            lines.append(f"{t.date.isoformat()} {prefix} {t.content}")
        if not lines:
            lines.append("（本周暂无记录）")
        progress_text = "\n".join(lines)
    add_page(progress_text, SECTION_TITLES["progress"])

    add_page(sections.get("pending", ""),
             "三、后续待解决的问题或研究方案调整：")
    add_page(sections.get("next_plan", ""),
             "四、下周详细计划：")
    add_page(sections.get("assessment", ""),
             "五、本人本周各课题进展自评（好、中、差）：")

    timeline_lines = []
    for m in milestones:
        planned = m.planned_date.isoformat() if m.planned_date else "（待定）"
        actual = m.actual_date.isoformat() if m.actual_date else "（未完成）"
        timeline_lines.append(
            f"{m.name}  |  计划：{planned}  |  实际：{actual}  |  {m.notes or ''}"
        )
    add_page("\n".join(timeline_lines),
             "六、本人成果发表时间节点更新：")

    pdf_path = os.path.join(OUTPUT_DIR, f"周报-方予之-{ws_str}.pdf")
    doc.save(pdf_path)
    doc.close()
    return pdf_path


def generate_weekly_report(week_start):
    """生成PPTX和PDF，返回 (pptx_path, pdf_path)"""
    pptx_path = generate_pptx(week_start)
    pdf_path = generate_pdf(week_start)
    return pptx_path, pdf_path
